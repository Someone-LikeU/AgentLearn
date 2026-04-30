#!/usr/bin/env python3
"""
本地文件语义索引构建器
扫描指定目录，提取文件内容，构建向量索引
"""

import os
import sys
import json
import argparse
import hashlib
from pathlib import Path
from datetime import datetime
from typing import List, Dict, Optional, Tuple

try:
    import chromadb
    from chromadb.config import Settings
    CHROMADB_AVAILABLE = True
except ImportError:
    CHROMADB_AVAILABLE = False

try:
    from sentence_transformers import SentenceTransformer
    SENTENCE_TRANSFORMERS_AVAILABLE = True
except ImportError:
    SENTENCE_TRANSFORMERS_AVAILABLE = False

try:
    import PyPDF2
    PYPDF2_AVAILABLE = True
except ImportError:
    PYPDF2_AVAILABLE = False

try:
    import docx
    PYTHON_DOCX_AVAILABLE = True
except ImportError:
    PYTHON_DOCX_AVAILABLE = False

try:
    from pptx import Presentation
    PYTHON_PPTX_AVAILABLE = True
except ImportError:
    PYTHON_PPTX_AVAILABLE = False


SKILL_DIR = Path(__file__).parent.parent
VECTOR_STORE_DIR = SKILL_DIR / "vector_store"
INDEX_METADATA_FILE = VECTOR_STORE_DIR / "metadata.json"
MODEL_CONFIG_FILE = SKILL_DIR / "model_config.json"

DEFAULT_FILE_TYPES = ['txt', 'md', 'pdf', 'docx', 'pptx']
CHUNK_SIZE = 500
CHUNK_OVERLAP = 50

DEFAULT_MODEL = 'moka-ai/m3e-base'
HUGGINGFACE_CACHE = Path.home() / '.cache' / 'huggingface' / 'hub'


def load_model_config() -> Dict:
    """加载模型配置"""
    if MODEL_CONFIG_FILE.exists():
        try:
            with open(MODEL_CONFIG_FILE, 'r', encoding='utf-8') as f:
                return json.load(f)
        except Exception:
            pass
    return {'model_name': DEFAULT_MODEL, 'local_path': None, 'use_local': True}


def find_local_model_path(model_name: str) -> Optional[str]:
    """查找本地模型路径"""
    cache_name = f"models--{model_name.replace('/', '--')}"
    cache_dir = HUGGINGFACE_CACHE / cache_name / 'snapshots'
    
    if cache_dir.exists():
        for snapshot in cache_dir.iterdir():
            if snapshot.is_dir():
                return str(snapshot)
    return None


def load_model_with_fallback() -> Tuple[Optional['SentenceTransformer'], str]:
    """加载模型（纯本地模式），失败时提供解决方案"""
    if not SENTENCE_TRANSFORMERS_AVAILABLE:
        return None, "sentence-transformers 未安装，请运行: pip install sentence-transformers"
    
    config = load_model_config()
    model_name = config.get('model_name', DEFAULT_MODEL)
    local_path = config.get('local_path')
    
    print(f"正在加载嵌入模型: {model_name}")
    
    # 纯本地模式：只从本地路径加载
    try:
        # 1. 优先使用配置中的本地路径
        if local_path and Path(local_path).exists():
            model = SentenceTransformer(local_path)
            return model, f"✓ 成功加载本地模型: {local_path}"
        
        # 2. 从HuggingFace缓存查找
        cache_path = find_local_model_path(model_name)
        if cache_path:
            model = SentenceTransformer(cache_path)
            return model, f"✓ 成功加载缓存模型: {cache_path}"
        
        # 3. 未找到模型
        error_msg = f"✗ 未找到本地模型: {model_name}\n"
        error_msg += "\n解决方案:\n"
        error_msg += "1. 扫描本地模型:\n"
        error_msg += "   python scripts/model_manager.py --scan\n"
        error_msg += "2. 交互式选择模型:\n"
        error_msg += "   python scripts/model_manager.py --interactive\n"
        error_msg += "3. 查看下载配置指南:\n"
        error_msg += "   python scripts/model_manager.py --guide\n"
        return None, error_msg
        
    except Exception as e:
        error_msg = f"✗ 模型加载失败: {e}\n"
        error_msg += "\n解决方案:\n"
        error_msg += "1. 扫描本地模型:\n"
        error_msg += "   python scripts/model_manager.py --scan\n"
        error_msg += "2. 交互式选择模型:\n"
        error_msg += "   python scripts/model_manager.py --interactive\n"
        error_msg += "3. 查看下载配置指南:\n"
        error_msg += "   python scripts/model_manager.py --guide\n"
        return None, error_msg


def get_file_hash(file_path: Path) -> str:
    """计算文件哈希值用于变更检测"""
    hasher = hashlib.md5()
    try:
        with open(file_path, 'rb') as f:
            for chunk in iter(lambda: f.read(8192), b''):
                hasher.update(chunk)
        return hasher.hexdigest()
    except Exception:
        return ""


def extract_text_from_txt(file_path: Path) -> List[str]:
    """提取 txt 文件内容"""
    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()
        return [content] if content.strip() else []
    except Exception as e:
        print(f"读取 txt 文件失败 {file_path}: {e}", file=sys.stderr)
        return []


def extract_text_from_md(file_path: Path) -> List[str]:
    """提取 markdown 文件内容"""
    return extract_text_from_txt(file_path)


def extract_text_from_pdf(file_path: Path) -> List[str]:
    """提取 PDF 文件内容"""
    if not PYPDF2_AVAILABLE:
        print("PyPDF2 未安装，跳过 PDF 文件", file=sys.stderr)
        return []
    
    chunks = []
    try:
        with open(file_path, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            for i, page in enumerate(reader.pages):
                text = page.extract_text()
                if text and text.strip():
                    chunks.append(f"[第{i+1}页] {text.strip()}")
    except Exception as e:
        print(f"读取 PDF 文件失败 {file_path}: {e}", file=sys.stderr)
    return chunks


def extract_text_from_docx(file_path: Path) -> List[str]:
    """提取 Word 文件内容"""
    if not PYTHON_DOCX_AVAILABLE:
        print("python-docx 未安装，跳过 docx 文件", file=sys.stderr)
        return []
    
    try:
        doc = docx.Document(str(file_path))
        content = '\n'.join([para.text for para in doc.paragraphs if para.text.strip()])
        return [content] if content else []
    except Exception as e:
        print(f"读取 docx 文件失败 {file_path}: {e}", file=sys.stderr)
        return []


def extract_text_from_pptx(file_path: Path) -> List[str]:
    """提取 PowerPoint 文件内容"""
    if not PYTHON_PPTX_AVAILABLE:
        print("python-pptx 未安装，跳过 pptx 文件", file=sys.stderr)
        return []
    
    chunks = []
    try:
        prs = Presentation(str(file_path))
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            if texts:
                chunks.append(f"[第{i+1}页] {' '.join(texts)}")
    except Exception as e:
        print(f"读取 pptx 文件失败 {file_path}: {e}", file=sys.stderr)
    return chunks


EXTRACTORS = {
    'txt': extract_text_from_txt,
    'md': extract_text_from_md,
    'pdf': extract_text_from_pdf,
    'docx': extract_text_from_docx,
    'pptx': extract_text_from_pptx,
}


def chunk_text(text: str, chunk_size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> List[str]:
    """将长文本分块"""
    if len(text) <= chunk_size:
        return [text]
    
    chunks = []
    start = 0
    while start < len(text):
        end = start + chunk_size
        chunk = text[start:end]
        chunks.append(chunk)
        start = end - overlap
    return chunks


def scan_directory(root_dir: Path, file_types: List[str], include_content: bool) -> List[Dict]:
    """扫描目录并提取文件信息"""
    files_data = []
    
    for ext in file_types:
        pattern = f'*.{ext}'
        for file_path in root_dir.rglob(pattern):
            if file_path.is_file():
                try:
                    stat = file_path.stat()
                    file_info = {
                        'path': str(file_path.absolute()),
                        'name': file_path.name,
                        'type': ext,
                        'size': stat.st_size,
                        'mtime': datetime.fromtimestamp(stat.st_mtime).isoformat() + 'Z',
                        'hash': get_file_hash(file_path),
                        'chunks': [],
                    }
                    
                    if include_content and ext in EXTRACTORS:
                        chunks = EXTRACTORS[ext](file_path)
                        file_info['chunks'] = chunks
                    
                    files_data.append(file_info)
                    print(f"已索引: {file_path.name}")
                    
                except Exception as e:
                    print(f"处理文件失败 {file_path}: {e}", file=sys.stderr)
    
    return files_data


def build_index(files_data: List[Dict], collection_name: str = "local_files"):
    """构建向量索引"""
    if not CHROMADB_AVAILABLE:
        print("chromadb 未安装，无法构建向量索引", file=sys.stderr)
        return None
    
    if not SENTENCE_TRANSFORMERS_AVAILABLE:
        print("sentence-transformers 未安装，无法构建向量索引", file=sys.stderr)
        return None
    
    model, message = load_model_with_fallback()
    if model is None:
        print(message, file=sys.stderr)
        return None
    
    print(message)
    
    print("正在初始化向量数据库...")
    VECTOR_STORE_DIR.mkdir(parents=True, exist_ok=True)
    client = chromadb.PersistentClient(path=str(VECTOR_STORE_DIR))
    
    try:
        collection = client.get_collection(name=collection_name)
        client.delete_collection(name=collection_name)
    except:
        pass
    
    collection = client.create_collection(
        name=collection_name,
        metadata={"hnsw:space": "cosine"}
    )
    
    print("正在构建索引...")
    ids = []
    documents = []
    metadatas = []
    
    for file_info in files_data:
        file_id = file_info['path']
        
        doc_text = f"文件名: {file_info['name']}\n"
        if file_info['chunks']:
            doc_text += f"内容摘要: {' '.join(file_info['chunks'][:3])[:500]}"
        
        ids.append(file_id)
        documents.append(doc_text)
        metadatas.append({
            'path': file_info['path'],
            'name': file_info['name'],
            'type': file_info['type'],
            'size': str(file_info['size']),
            'mtime': file_info['mtime'],
        })
    
    if documents:
        batch_size = 100
        for i in range(0, len(documents), batch_size):
            batch_ids = ids[i:i+batch_size]
            batch_docs = documents[i:i+batch_size]
            batch_meta = metadatas[i:i+batch_size]
            
            embeddings = model.encode(batch_docs, show_progress_bar=False).tolist()
            collection.add(
                ids=batch_ids,
                documents=batch_docs,
                embeddings=embeddings,
                metadatas=batch_meta
            )
    
    print(f"索引构建完成，共 {len(documents)} 个文档")
    return collection


def save_metadata(files_data: List[Dict], root_dir: Path, file_types: List[str], include_content: bool):
    """保存索引元数据"""
    metadata = {
        'root_dir': str(root_dir.absolute()),
        'file_types': file_types,
        'include_content': include_content,
        'total_files': len(files_data),
        'indexed_at': datetime.now().isoformat(),
        'files': files_data
    }
    
    with open(INDEX_METADATA_FILE, 'w', encoding='utf-8') as f:
        json.dump(metadata, f, ensure_ascii=False, indent=2)
    
    print(f"元数据已保存到 {INDEX_METADATA_FILE}")


def main():
    parser = argparse.ArgumentParser(description='本地文件语义索引构建器')
    parser.add_argument('--dir', type=str, default='.', help='搜索根目录')
    parser.add_argument('--include-content', action='store_true', default=True, help='是否索引文件内容')
    parser.add_argument('--file-types', type=str, default='txt,md,pdf,docx,pptx', help='文件类型，逗号分隔')
    parser.add_argument('--force', action='store_true', help='强制重建索引（全量刷新）')
    parser.add_argument('--incremental', action='store_true', help='增量更新，只处理变化的文件')
    
    args = parser.parse_args()
    
    root_dir = Path(args.dir).resolve()
    file_types = [ft.strip() for ft in args.file_types.split(',')]
    
    print(f"搜索目录: {root_dir}")
    print(f"文件类型: {file_types}")
    print(f"包含内容: {args.include_content}")
    
    if args.force and args.incremental:
        print("警告: --force 和 --incremental 不能同时使用，将使用全量刷新模式", file=sys.stderr)
        args.incremental = False
    
    if not root_dir.exists():
        print(f"目录不存在: {root_dir}", file=sys.stderr)
        sys.exit(1)
    
    if args.incremental:
        print("\n增量更新模式...")
        existing_metadata = load_existing_metadata()
        
        if not existing_metadata or Path(existing_metadata.get('root_dir', '')).resolve() != root_dir:
            print("无现有索引或目录不匹配，将执行全量构建")
            args.incremental = False
        else:
            existing_files = {f['path']: f['hash'] for f in existing_metadata.get('files', [])}
            print(f"现有索引文件数: {len(existing_files)}")
    
    print("\n开始扫描文件...")
    files_data = scan_directory(root_dir, file_types, args.include_content)
    
    if not files_data:
        print("未找到匹配的文件")
        return
    
    print(f"\n找到 {len(files_data)} 个文件")
    
    if args.incremental and existing_metadata:
        new_files = []
        updated_files = []
        deleted_files = []
        
        current_files = {f['path']: f for f in files_data}
        
        for path, hash_val in existing_files.items():
            if path not in current_files:
                deleted_files.append(path)
        
        for file_info in files_data:
            path = file_info['path']
            if path not in existing_files:
                new_files.append(file_info)
            elif file_info['hash'] != existing_files[path]:
                updated_files.append(file_info)
        
        print(f"新增文件: {len(new_files)}")
        print(f"更新文件: {len(updated_files)}")
        print(f"删除文件: {len(deleted_files)}")
        
        files_to_process = new_files + updated_files
        
        if files_to_process:
            print("\n开始增量更新向量索引...")
            incremental_update(files_to_process, deleted_files)
        else:
            print("\n无文件变更，索引已是最新")
        
        all_files = [f for f in files_data if f['path'] not in deleted_files]
        save_metadata(all_files, root_dir, file_types, args.include_content)
    else:
        print("\n开始构建向量索引...")
        build_index(files_data)
        save_metadata(files_data, root_dir, file_types, args.include_content)
    
    print("\n索引构建完成！")


def load_existing_metadata() -> Optional[Dict]:
    """加载现有索引元数据"""
    if not INDEX_METADATA_FILE.exists():
        return None
    try:
        with open(INDEX_METADATA_FILE, 'r', encoding='utf-8') as f:
            return json.load(f)
    except Exception:
        return None


def incremental_update(new_files: List[Dict], deleted_paths: List[str]):
    """增量更新向量索引"""
    if not CHROMADB_AVAILABLE or not SENTENCE_TRANSFORMERS_AVAILABLE:
        print("依赖未安装，无法执行增量更新", file=sys.stderr)
        return
    
    model, message = load_model_with_fallback()
    if model is None:
        print(message, file=sys.stderr)
        return
    
    print(message)
    
    print("正在连接向量数据库...")
    VECTOR_STORE_DIR.mkdir(parents=True, exist_ok=True)
    client = chromadb.PersistentClient(path=str(VECTOR_STORE_DIR))
    
    try:
        collection = client.get_collection(name="local_files")
    except:
        collection = client.create_collection(
            name="local_files",
            metadata={"hnsw:space": "cosine"}
        )
    
    if deleted_paths:
        print(f"正在删除 {len(deleted_paths)} 个文件的索引...")
        try:
            collection.delete(ids=deleted_paths)
        except Exception as e:
            print(f"删除失败: {e}", file=sys.stderr)
    
    if new_files:
        print(f"正在添加 {len(new_files)} 个新/更新文件...")
        ids = []
        documents = []
        metadatas = []
        
        for file_info in new_files:
            doc_text = f"文件名: {file_info['name']}\n"
            if file_info['chunks']:
                doc_text += f"内容摘要: {' '.join(file_info['chunks'][:3])[:500]}"
            
            ids.append(file_info['path'])
            documents.append(doc_text)
            metadatas.append({
                'path': file_info['path'],
                'name': file_info['name'],
                'type': file_info['type'],
                'size': str(file_info['size']),
                'mtime': file_info['mtime'],
            })
        
        batch_size = 100
        for i in range(0, len(documents), batch_size):
            batch_ids = ids[i:i+batch_size]
            batch_docs = documents[i:i+batch_size]
            batch_meta = metadatas[i:i+batch_size]
            
            embeddings = model.encode(batch_docs, show_progress_bar=False).tolist()
            
            try:
                collection.upsert(
                    ids=batch_ids,
                    documents=batch_docs,
                    embeddings=embeddings,
                    metadatas=batch_meta
                )
            except Exception as e:
                print(f"添加失败: {e}", file=sys.stderr)
    
    print(f"增量更新完成")


if __name__ == '__main__':
    main()
